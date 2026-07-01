# flashcard_mcq_app.py
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, simpledialog
import random
import textwrap
import os
import pathlib

# NLTK + WordNet  this is my current code add to this
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords, wordnet

# PDF export (reportlab)
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

# ---------- Optional file parsers (loaded if installed) ----------
try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

# .doc (legacy) via textract (optional; heavy dependency)
_textract_error = None
try:
    import textract  # noqa
except Exception as _e:
    textract = None
    _textract_error = _e

# -------------------- Ensure NLTK data --------------------
# These downloads are idempotent and quiet; they will not re-download if present.
nltk_packages = ['punkt', 'averaged_perceptron_tagger', 'wordnet', 'omw-1.4', 'stopwords']
for pkg in nltk_packages:
    try:
        nltk.data.find(pkg if pkg != 'omw-1.4' else 'corpora/omw-1.4')
    except LookupError:
        nltk.download(pkg, quiet=True)

STOPWORDS = set(stopwords.words('english'))

# -------------------- File text extraction helpers --------------------
def extract_text_from_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_text_from_pdf(path):
    if PyPDF2 is None:
        raise ImportError("Missing dependency: PyPDF2 (pip install PyPDF2)")
    text_parts = []
    with open(path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            try:
                text_parts.append(page.extract_text() or "")
            except Exception:
                # Skip problematic pages but keep going
                pass
    return "\n".join(text_parts)

def extract_text_from_docx(path):
    if DocxDocument is None:
        raise ImportError("Missing dependency: python-docx (pip install python-docx)")
    doc = DocxDocument(path)
    chunks = []
    # paragraphs
    for p in doc.paragraphs:
        if p.text:
            chunks.append(p.text)
    # tables
    for t in doc.tables:
        for row in t.rows:
            row_text = [cell.text for cell in row.cells if cell.text]
            if row_text:
                chunks.append(" | ".join(row_text))
    return "\n".join(chunks)

def extract_text_from_pptx(path):
    if Presentation is None:
        raise ImportError("Missing dependency: python-pptx (pip install python-pptx)")
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)

def extract_text_from_doc(path):
    # Legacy .doc requires textract or other external tools
    if textract is None:
        raise ImportError("Missing dependency for .doc: textract (pip install textract)\n"
                          f"Loader error: {_textract_error}")
    data = textract.process(path)
    return data.decode(errors="ignore")

def extract_text_smart(path):
    ext = pathlib.Path(path).suffix.lower()
    if ext == ".txt":
        return extract_text_from_txt(path)
    elif ext == ".pdf":
        return extract_text_from_pdf(path)
    elif ext == ".docx":
        return extract_text_from_docx(path)
    elif ext == ".pptx":
        return extract_text_from_pptx(path)
    elif ext == ".doc":
        return extract_text_from_doc(path)
    else:
        # Fallback: try as text
        return extract_text_from_txt(path)

# -------------------- MCQ generation utilities --------------------
def get_wordnet_distractors(word):
    """Return a list of candidate distractors from WordNet for `word`."""
    word = word.strip()
    distractors = set()
    synsets = wordnet.synsets(word)
    if not synsets:
        return []
    # Collect synonyms/related lemmas and hyponyms
    for syn in synsets:
        # synonyms (lemmas)
        for lemma in syn.lemmas():
            name = lemma.name().replace('_', ' ')
            if name.lower() != word.lower() and name.isalpha():
                distractors.add(name)
        # hyponyms (more specific)
        for hy in syn.hyponyms():
            for lemma in hy.lemmas():
                name = lemma.name().replace('_', ' ')
                if name.lower() != word.lower() and name.isalpha():
                    distractors.add(name)
        # hypernyms -> other hyponyms (siblings)
        for hyper in syn.hypernyms():
            for hypo in hyper.hyponyms():
                for lemma in hypo.lemmas():
                    name = lemma.name().replace('_', ' ')
                    if name.lower() != word.lower() and name.isalpha():
                        distractors.add(name)
    # filtering & return list
    distractors = [d for d in distractors if d.lower() != word.lower()]
    # dedupe case-insensitively while preserving readability
    seen = set()
    out = []
    for d in distractors:
        key = d.lower()
        if key not in seen:
            seen.add(key)
            out.append(d)
    return out

def choose_keyword_for_sentence(sentence):
    """Select a suitable noun/proper noun token from a sentence as the keyword."""
    tokens = word_tokenize(sentence)
    tagged = nltk.pos_tag(tokens)
    candidates = []
    for word, pos in tagged:
        # prefer proper nouns and nouns longer than 2 chars, alphabetic and not a stopword
        if pos in ('NN', 'NNS', 'NNP', 'NNPS') and word.isalpha() and word.lower() not in STOPWORDS and len(word) > 2:
            candidates.append(word)
    if not candidates:
        return None
    # prefer longer words (more distinctive)
    candidates.sort(key=lambda w: len(w), reverse=True)
    return candidates[0]

def generate_mcqs_from_text(text, num_questions=5):
    """
    Produce up to num_questions MCQs strictly from given text.
    Each MCQ: {'question': <sentence with blank>, 'options': [..4..], 'answer': <word>}
    Uses WordNet distractors only; if insufficient distractors, the sentence is skipped.
    """
    sentences = sent_tokenize(text)
    mcqs = []
    # iterate sentences in order; try to extract one mcq per sentence
    for sent in sentences:
        if len(mcqs) >= num_questions:
            break
        kw = choose_keyword_for_sentence(sent)
        if not kw:
            continue
        distractors = get_wordnet_distractors(kw)
        # Need at least 3 distractors
        if len(distractors) < 3:
            continue
        # choose 3 distractors
        opts = random.sample(distractors, 3)
        opts.append(kw)
        random.shuffle(opts)
        q_text = sent.replace(kw, "_____")
        mcqs.append({
            'question': q_text,
            'options': opts,
            'answer': kw
        })
    return mcqs

# -------------------- GUI App --------------------
class FlashcardMCQApp:
    def __init__(self, root):
        self.root = root
        root.title("AI Flashcard MCQ Maker")
        root.geometry("920x620")
        root.minsize(780, 520)

        # App state
        self.content_text = ""
        self.mcqs = []
        self.num_mcqs = 5

        # Quiz runtime
        self.deck = []         # shuffled mcqs for play
        self.index = 0
        self.score = 0
        self.user_answers = []

        # Colors & style
        self.card_colors = ["#FFCDD2", "#F8BBD0", "#E1BEE7", "#BBDEFB", "#B2DFDB", "#C8E6C9", "#FFF9C4", "#FFE0B2"]
        self.right_color = "#A5D6A7"
        self.wrong_color = "#EF9A9A"
        self.bg_color = "#F5F7FA"

        # Build UI (single window with frames)
        self.build_ui()

    # -------- UI Construction --------
    def build_ui(self):
        self.root.configure(bg=self.bg_color)
        # Top frame - title
        top = tk.Frame(self.root, bg=self.bg_color)
        top.pack(fill='x', padx=12, pady=(12,0))
        title = tk.Label(top, text="AI Flashcard MCQ Maker", font=("Helvetica", 20, "bold"), bg=self.bg_color)
        title.pack(side='left', padx=6)

        subtitle = tk.Label(top, text="— NLTK + WordNet MCQs (content-only)", font=("Helvetica", 10), bg=self.bg_color, fg="#333")
        subtitle.pack(side='left', padx=8)

        # Main frames
        main = tk.Frame(self.root, bg=self.bg_color)
        main.pack(fill='both', expand=True, padx=12, pady=12)

        # Left: Input & options
        left = tk.Frame(main, width=360, bg=self.bg_color)
        left.pack(side='left', fill='y', padx=(0,12))

        self.build_input_panel(left)

        # Right: Preview / Quiz area
        right = tk.Frame(main, bg=self.bg_color)
        right.pack(side='right', fill='both', expand=True)

        self.build_quiz_panel(right)

    def build_input_panel(self, parent):
        # Input label
        lbl = tk.Label(parent, text="1) Paste or Upload Content", font=("Helvetica", 12, "bold"), bg=self.bg_color)
        lbl.pack(anchor='w', pady=(2,6))

        self.textbox = tk.Text(parent, height=14, wrap='word', font=("Arial", 11))
        self.textbox.pack(fill='x')

        btn_row = tk.Frame(parent, bg=self.bg_color)
        btn_row.pack(fill='x', pady=8)
        tk.Button(btn_row, text="Open File (.txt/.pdf/.docx/.pptx/.doc)", command=self.load_file, bg="#1976D2", fg="white").pack(side='left', padx=4)
        tk.Button(btn_row, text="Clear Text", command=lambda: self.textbox.delete("1.0", tk.END)).pack(side='left', padx=4)

        # Quiz customization
        tk.Label(parent, text="2) Customization", font=("Helvetica", 12, "bold"), bg=self.bg_color).pack(anchor='w', pady=(12,6))

        opts_frame = tk.Frame(parent, bg=self.bg_color)
        opts_frame.pack(fill='x')
        tk.Label(opts_frame, text="Number of MCQs:", bg=self.bg_color).grid(row=0, column=0, sticky='w')
        self.spin_mcq = tk.Spinbox(opts_frame, from_=0, to=50, width=6)
        self.spin_mcq.grid(row=0, column=1, sticky='w', padx=8)

        # Control buttons
        control = tk.Frame(parent, bg=self.bg_color)
        control.pack(fill='x', pady=(14,0))
        tk.Button(control, text="Generate MCQs", command=self.generate_mcqs_btn, bg="#388E3C", fg="white").pack(fill='x', pady=6)
        tk.Button(control, text="Play Quiz (Flashcards)", command=self.play_quiz_btn, bg="#F57C00", fg="white").pack(fill='x', pady=6)
        tk.Button(control, text="Export PDF (Questions + Answers)", command=self.export_pdf, bg="#455A64", fg="white").pack(fill='x', pady=6)

        # status
        self.status_var = tk.StringVar(value="Ready")
        tk.Label(parent, textvariable=self.status_var, bg=self.bg_color, fg="#333").pack(anchor='w', pady=(10,0))

    def build_quiz_panel(self, parent):
        # Card area
        self.card_frame = tk.Frame(parent, bg=self.bg_color)
        self.card_frame.pack(fill='both', expand=True)

        # placeholder message
        self.placeholder = tk.Label(self.card_frame,
            text="Generate MCQs from your content,\nor press Play Quiz to begin.",
            font=("Helvetica", 16), bg=self.bg_color, fg="#666", justify='center')
        self.placeholder.pack(expand=True)

        # card (hidden until quiz starts)
        self.card = tk.Frame(self.card_frame, width=520, height=360, bg="white", relief='raised', bd=2)
        # center geometry management
        self.card.place_forget()

        # question label
        self.q_label = tk.Label(self.card, text="", wraplength=480, font=("Helvetica", 16, "bold"), bg="white", justify='center')
        self.q_label.pack(pady=(26,8), padx=12)

        # option buttons container
        self.option_buttons = []
        opts_cont = tk.Frame(self.card, bg="white")
        opts_cont.pack(pady=8)
        for i in range(4):
            b = tk.Button(opts_cont, text="", font=("Helvetica", 13), width=34, relief='raised',
                          command=lambda idx=i: self.select_option(idx))
            b.pack(pady=6)
            self.option_buttons.append(b)

        # bottom controls on card
        bottom = tk.Frame(self.card, bg="white")
        bottom.pack(side='bottom', pady=10, fill='x')
        self.progress_label = tk.Label(bottom, text="", bg="white")
        self.progress_label.pack(side='left', padx=10)
        self.next_btn = tk.Button(bottom, text="Skip →", command=self.skip_question, bg="#1976D2", fg="white")
        self.next_btn.pack(side='right', padx=12)

        # Score meter frame (shown after quiz)
        self.score_frame = tk.Frame(self.card_frame, bg=self.bg_color)
        # includes progressbar and message
        self.score_msg = tk.Label(self.score_frame, text="", font=("Helvetica", 18, "bold"), bg=self.bg_color)
        self.score_msg.pack(pady=20)
        self.score_progress = ttk.Progressbar(self.score_frame, orient='horizontal', length=400, mode='determinate')
        self.score_progress.pack(pady=8)

    # ---------- Handlers ----------
    def load_file(self):
        path = filedialog.askopenfilename(
            filetypes=[
                ("Supported files", "*.txt *.pdf *.docx *.pptx *.doc"),
                ("Text files", "*.txt"),
                ("PDF files", "*.pdf"),
                ("Word (DOCX)", "*.docx"),
                ("PowerPoint (PPTX)", "*.pptx"),
                ("Word (DOC)", "*.doc"),
                ("All files", "*.*"),
            ]
        )
        if not path:
            return
        try:
            data = extract_text_smart(path)
            if not data.strip():
                messagebox.showwarning("No text found", "Could not extract text from this file.")
                return
            self.textbox.delete("1.0", tk.END)
            self.textbox.insert(tk.END, data.strip())
            self.status_var.set(f"Loaded file: {os.path.basename(path)}")
        except ImportError as ie:
            messagebox.showerror("Missing dependency", str(ie))
        except Exception as e:
            messagebox.showerror("Error", f"Could not read file: {e}")

    def generate_mcqs_btn(self):
        text = self.textbox.get("1.0", tk.END).strip()
        if not text:
            messagebox.showerror("Error", "Please paste or load content first.")
            return
        try:
            n = int(self.spin_mcq.get())
            if n < 0:
                raise ValueError
        except Exception:
            messagebox.showerror("Error", "Enter a valid number of MCQs (0 or positive).")
            return

        self.status_var.set("Generating MCQs... (this may take a few seconds)")
        self.root_update()
        mcqs = generate_mcqs_from_text(text, num_questions=n)
        if not mcqs:
            messagebox.showinfo("No MCQs", "Could not generate MCQs from this content with WordNet distractors.\nTry a different text or a smaller number.")
            self.status_var.set("Ready")
            return
        self.mcqs = mcqs
        self.num_mcqs = len(mcqs)
        self.status_var.set(f"Generated {len(mcqs)} MCQs. You can Play Quiz or Export PDF.")
        messagebox.showinfo("Done", f"Generated {len(mcqs)} MCQs from the content.")

    def play_quiz_btn(self):
        # If generation not done, generate automatically using chosen number
        if not self.mcqs:
            # attempt auto-generate using spin value
            self.generate_mcqs_btn()
            if not self.mcqs:
                return

        # create a shuffled deck
        self.deck = list(self.mcqs)
        random.shuffle(self.deck)
        self.index = 0
        self.score = 0
        self.user_answers = []
        self.show_card_for_current()

    def show_card_for_current(self):
        if self.index >= len(self.deck):
            self.show_score()
            return

        q = self.deck[self.index]
        # show card center
        self.placeholder.pack_forget()
        self.score_frame.pack_forget()
        self.card.place(relx=0.5, rely=0.5, anchor='center')
        # set a colorful bg
        color = random.choice(self.card_colors)
        self.card.configure(bg=color)
        self.q_label.configure(text=f"Q{self.index+1}: {q['question']}", bg=color)
        for b, opt in zip(self.option_buttons, q['options']):
            b.configure(text=opt, bg="white", state='normal', relief='raised')
        self.progress_label.configure(text=f"{self.index+1}/{len(self.deck)}", bg=color)
        self.next_btn.configure(text="Skip →", bg="#1976D2", fg="white", state='normal')

    def select_option(self, opt_idx):
        """User clicked an option by index."""
        # disable options to prevent double clicks
        for b in self.option_buttons:
            b.config(state='disabled')
        selected_text = self.option_buttons[opt_idx]['text']
        correct = self.deck[self.index]['answer']
        # color feedback
        if selected_text == correct:
            self.option_buttons[opt_idx].config(bg=self.right_color)
            self.score += 1
        else:
            self.option_buttons[opt_idx].config(bg=self.wrong_color)
            # highlight correct option
            for b in self.option_buttons:
                if b['text'] == correct:
                    b.config(bg=self.right_color)
                    break
        self.user_answers.append(selected_text)
        # small pause then next question
        self.root.after(700, self.advance_card)

    def skip_question(self):
        # record blank or "SKIPPED"
        self.user_answers.append("SKIPPED")
        # briefly show the correct answer then advance
        correct = self.deck[self.index]['answer']
        for b in self.option_buttons:
            if b['text'] == correct:
                b.config(bg=self.right_color)
            b.config(state='disabled')
        self.root.after(700, self.advance_card)

    def advance_card(self):
        # reset button colors
        for b in self.option_buttons:
            b.config(bg='white', state='normal', relief='raised')
        self.index += 1
        if self.index < len(self.deck):
            self.show_card_for_current()
        else:
            self.show_score()

    def show_score(self):
        # hide card
        self.card.place_forget()
        # show a score frame with progressbar
        self.score_msg.configure(text=f"Your Score: {self.score} / {len(self.deck)}")
        self.score_progress['maximum'] = len(self.deck)
        self.score_progress['value'] = 0
        self.score_frame.pack(expand=True)
        # animate progress bar up to score
        self.animate_score(0)

    def animate_score(self, v):
        if v > self.score:
            return
        self.score_progress['value'] = v
        # message / feedback
        percent = (self.score / len(self.deck)) * 100 if len(self.deck) > 0 else 0
        if percent >= 80:
            msg = "🎉 Excellent!"
        elif percent >= 50:
            msg = "🙂 Good job!"
        else:
            msg = "📚 Keep practicing!"
        self.score_msg.configure(text=f"Your Score: {self.score} / {len(self.deck)}\n{msg}")
        if v < self.score:
            self.root.after(80, lambda: self.animate_score(v+1))
        else:
            # show buttons after animation
            btns = tk.Frame(self.score_frame, bg=self.bg_color)
            btns.pack(pady=12)
            tk.Button(btns, text="Play Again", command=self.reset_to_start, bg="#1976D2", fg="white").pack(side='left', padx=6)
            tk.Button(btns, text="Export PDF (clean)", command=self.export_pdf, bg="#455A64", fg="white").pack(side='left', padx=6)

    def reset_to_start(self):
        self.score_frame.pack_forget()
        self.card.place_forget()
        self.placeholder.pack(expand=True)
        self.mcqs = []
        self.deck = []
        self.index = 0
        self.score = 0
        self.user_answers = []
        self.status_var.set("Ready")

    # ---------- PDF Export ----------
    def export_pdf(self):
        # if mcqs empty ask to generate
        if not self.mcqs:
            # try to auto generate once
            text = self.textbox.get("1.0", tk.END).strip()
            if not text:
                messagebox.showerror("No content", "Paste or upload content and generate MCQs first.")
                return
            try:
                n = int(self.spin_mcq.get())
                if n < 0:
                    raise ValueError
            except Exception:
                messagebox.showerror("Error", "Enter a valid number of MCQs (0 or positive).")
                return
            self.mcqs = generate_mcqs_from_text(text, num_questions=n)
            if not self.mcqs:
                messagebox.showinfo("No MCQs", "Could not generate MCQs from content.")
                return

        save_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files","*.pdf")], title="Save PDF")
        if not save_path:
            return

        try:
            c = canvas.Canvas(save_path, pagesize=A4)
            width, height = A4
            margin = 50
            y = height - margin
            c.setFont("Helvetica-Bold", 18)
            c.drawString(margin, y, "Flashcard Quiz — Questions, Options & Answers")
            y -= 30
            c.setFont("Helvetica", 12)

            for i, q in enumerate(self.mcqs, start=1):
                # Question
                qtxt = f"Q{i}. {q['question']}"
                wrapped_q = textwrap.wrap(qtxt, 90)
                for line in wrapped_q:
                    c.drawString(margin, y, line)
                    y -= 16
                    if y < margin:
                        c.showPage(); y = height - margin

                # Options
                for idx, opt in enumerate(q['options'], start=1):
                    opt_line = f"   {chr(64+idx)}. {opt}"
                    wrapped_opt = textwrap.wrap(opt_line, 85)
                    for line in wrapped_opt:
                        c.drawString(margin + 10, y, line)
                        y -= 14
                        if y < margin:
                            c.showPage(); y = height - margin

                # Correct answer
                ans_line = f"   Correct Answer: {q['answer']}"
                c.setFont("Helvetica-Oblique", 11)
                c.drawString(margin + 10, y, ans_line)
                c.setFont("Helvetica", 12)
                y -= 20
                if y < margin:
                    c.showPage(); y = height - margin

            c.save()
            messagebox.showinfo("Saved", f"PDF saved to {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to save PDF: {e}")

    def root_update(self):
        self.root.update_idletasks()
        self.root.update()

# -------------------- Run App --------------------
if __name__ == "__main__":
    root = tk.Tk()
    app = FlashcardMCQApp(root)
    root.mainloop()